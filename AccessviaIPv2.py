import os
import socket
import requests
from openpyxl import Workbook
from selenium import webdriver
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def check_accessibility(url):
    try:
        response = requests.get(url, verify=False)  # Disable SSL certificate verification
        return response.status_code == 200
    except Exception as e:
        print(f"Error accessing {url}: {e}")
        return False

def check_access_via_ip(domain):
    try:
        ip = socket.gethostbyname(domain)
        return True, ip
    except socket.error:
        return False, None

def take_screenshot(ip, output_folder, idx):
    try:
        options = webdriver.FirefoxOptions()
        options.headless = True
        driver = webdriver.Firefox(options=options)
        driver.get(f"http://{ip}")  # Access the IP address
        screenshot_path = os.path.join(output_folder, f"{idx}_{ip.replace('.', '_')}.png")
        driver.save_screenshot(screenshot_path)
        driver.quit()
        return screenshot_path
    except Exception as e:
        print(f"Failed to capture screenshot of {ip}: {e}")
        return None

def get_results_from_file(file_path, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    results = []
    with open(file_path, 'r') as file:
        for idx, line in enumerate(file, start=1):
            domain = line.strip()
            access_via_ip, ip = check_access_via_ip(domain)
            if access_via_ip:
                screenshot_path = take_screenshot(ip, output_folder, idx)
                accessible = check_accessibility(f"http://{ip}")
                results.append((idx, domain, access_via_ip, ip, accessible, screenshot_path))
            else:
                results.append((idx, domain, access_via_ip, ip, False, None))
    return results

def save_to_excel(results, output_excel):
    wb = Workbook()
    ws = wb.active
    ws.title = "Website Results"
    headers = ["#", "Website", "Access via IP", "IP Address", "Accessible via IP", "Screenshot Path"]
    ws.append(headers)

    for row in results:
        ws.append(row)

    # Apply formatting
    for col in ws.columns:
        max_length = 0
        column = col[0].column_letter
        for cell in col:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(cell.value)
            except:
                pass
        adjusted_width = (max_length + 2) * 1.2
        ws.column_dimensions[column].width = adjusted_width

    wb.save(output_excel + ".xlsx")  # Save as Excel file with ".xlsx" extension
    print(f"Results saved to {output_excel}.xlsx")

def save_to_pptx(results, output_pptx):
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[5]

    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Website Screenshots"
    subtitle = slide.placeholders[1]
    subtitle.text = "Checker for IP address access"

    # Add content slides with screenshots
    for idx, _, _, _, _, screenshot_path in results:
        if screenshot_path:
            slide = prs.slides.add_slide(content_slide_layout)
            left = top = Inches(1)
            pic = slide.shapes.add_picture(screenshot_path, left, top, width=Inches(8))
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
            text_frame = title_shape.text_frame
            p = text_frame.add_paragraph()
            p.text = f"Screenshot {idx}:"
            p.font.bold = True
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.LEFT

    prs.save(output_pptx + ".pptx")
    print(f"Screenshots saved to {output_pptx}.pptx")

if __name__ == "__main__":
    file_path = input("Enter the path to the file containing URLs: ").strip()
    output_excel = input("Enter the name of the output Excel file: ").strip()
    output_pptx = input("Enter the name of the output PowerPoint file: ").strip()
    output_folder = input("Enter the folder to save screenshots: ").strip()

    results = get_results_from_file(file_path, output_folder)
    save_to_excel(results, output_excel)
    save_to_pptx(results, output_pptx)
